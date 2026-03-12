#!/usr/bin/env python3
"""Generate FFA-Net literature-sharing PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy, os

# ── Palette ──────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x0C, 0x29)
BG_BLUE   = RGBColor(0x0D, 0x1B, 0x2A)
BG_TEAL   = RGBColor(0x0B, 0x3D, 0x2E)
BG_BLACK  = RGBColor(0x0A, 0x0A, 0x0A)
BG_PURPLE = RGBColor(0x1A, 0x05, 0x33)

GOLD   = RGBColor(0xFF, 0xD2, 0x00)
CYAN   = RGBColor(0x00, 0xC6, 0xFF)
GREEN  = RGBColor(0x43, 0xE9, 0x7B)
ORANGE = RGBColor(0xF7, 0x97, 0x1E)
RED    = RGBColor(0xFF, 0x6B, 0x6B)
PURPLE = RGBColor(0xBF, 0x5A, 0xF2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0xAA, 0xAA, 0xAA)
DGRAY  = RGBColor(0x66, 0x66, 0x66)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helper utilities ──────────────────────────────────────

def add_slide(bg_color=BG_DARK):
    slide = prs.slides.add_slide(blank_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide

def txb(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def section_label(slide, text, color=CYAN):
    r = rect(slide, Inches(0.5), Inches(0.35), Inches(2.6), Inches(0.32),
             fill_color=RGBColor(0x22,0x22,0x44), line_color=None)
    txb(slide, f"◆  {text}", Inches(0.52), Inches(0.35), Inches(2.6), Inches(0.32),
        size=10, color=color)

def bottom_bar(slide, colors=None):
    """Gradient-ish accent bar at bottom."""
    segs = [ORANGE, GOLD, CYAN, GREEN]
    seg_w = W / len(segs)
    for i, c in enumerate(segs):
        rect(slide, i*seg_w, H - Inches(0.07), seg_w, Inches(0.07), fill_color=c)

def slide_num(slide, n, total=13):
    txb(slide, f"{n} / {total}",
        W - Inches(1.0), H - Inches(0.38), Inches(0.9), Inches(0.3),
        size=10, color=DGRAY, align=PP_ALIGN.RIGHT)

def divider(slide, y, color=DGRAY):
    r = rect(slide, Inches(0.5), y, W - Inches(1.0), Inches(0.015), fill_color=color)

# ════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ════════════════════════════════════════════════════════
s = add_slide(BG_DARK)

# background decorative circles (faked with rounded rect)
deco = s.shapes.add_shape(9, Inches(9.8), Inches(0.3), Inches(2.8), Inches(2.8))  # oval
deco.fill.solid(); deco.fill.fore_color.rgb = RGBColor(0x28,0x18,0x55)
deco.line.fill.background()

# badge
rect(s, Inches(0.5), Inches(0.7), Inches(3.2), Inches(0.38),
     fill_color=RGBColor(0x25,0x1a,0x4a))
txb(s, "AAAI 2020  ·  Official Implementation",
    Inches(0.52), Inches(0.72), Inches(3.2), Inches(0.36),
    size=11, color=GOLD, bold=True)

txb(s, "FFA-Net", Inches(0.5), Inches(1.25), Inches(9), Inches(1.1),
    size=60, bold=True, color=WHITE)
txb(s, "Feature Fusion Attention Network", Inches(0.5), Inches(2.25), Inches(9), Inches(0.55),
    size=26, color=RGBColor(0xBB,0xBB,0xCC))
txb(s, "for Single Image Dehazing", Inches(0.5), Inches(2.75), Inches(9), Inches(0.55),
    size=26, color=RGBColor(0xBB,0xBB,0xCC))

txb(s, "arXiv: 1911.07559", Inches(0.5), Inches(3.35), Inches(5), Inches(0.4),
    size=12, color=DGRAY)

divider(s, Inches(3.9), RGBColor(0x44,0x33,0x77))

txb(s, "Xu Qin, Zhilin Wang, Yuanchao Bai, Xiaodong Xie, Huizhu Jia",
    Inches(0.5), Inches(4.05), Inches(10), Inches(0.42), size=13, color=GRAY)
txb(s, "北京大学  ·  北京航空航天大学",
    Inches(0.5), Inches(4.5), Inches(7), Inches(0.42), size=13, color=GRAY)

# tags
for i,(txt,clr,bgc) in enumerate([
    ("AAAI 2020", GOLD, RGBColor(0x40,0x35,0x00)),
    ("图像复原",  CYAN, RGBColor(0x00,0x25,0x3a)),
    ("注意力机制",GREEN,RGBColor(0x05,0x2a,0x18)),
]):
    rx = Inches(0.5) + i * Inches(1.85)
    rect(s, rx, Inches(5.1), Inches(1.7), Inches(0.42), fill_color=bgc)
    txb(s, txt, rx, Inches(5.12), Inches(1.7), Inches(0.4), size=12, color=clr,
        bold=True, align=PP_ALIGN.CENTER)

bottom_bar(s)
slide_num(s, 1)

# ════════════════════════════════════════════════════════
# SLIDE 2 – Background
# ════════════════════════════════════════════════════════
s = add_slide(RGBColor(0x10,0x14,0x28))
section_label(s, "研究背景", CYAN)
txb(s, "为什么研究图像去雾？", Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=CYAN)

cards = [
    ("🌫️", "问题来源",   "雾霾导致光线散射，图像出现颜色失真、对比度降低、细节模糊等问题。"),
    ("🚗", "应用价值",   "自动驾驶感知、城市监控、卫星遥感等下游任务均需清晰图像输入。"),
    ("📐", "物理模型",   "大气散射模型：I(x) = J(x)·t(x) + A·(1-t(x))\n其中 t 为透射率，A 为大气光。"),
    ("⚡", "核心挑战",   "t(x) 与 A 均为未知，问题高度不适定；不同浓度、场景差异极大。"),
    ("🧠", "深度学习方案","端到端学习直接预测清晰图像，避免手工估算透射率，性能大幅超越传统方法。"),
    ("📊", "评估基准",   "RESIDE 数据集提供室内/室外场景，以 PSNR & SSIM 为量化评估指标。"),
]
cols, rows = 3, 2
cw, ch = Inches(3.9), Inches(1.8)
ox, oy = Inches(0.5), Inches(1.55)
for i,(icon,title,desc) in enumerate(cards):
    c,r = i%cols, i//cols
    x = ox + c*(cw+Inches(0.2))
    y = oy + r*(ch+Inches(0.18))
    rect(s, x, y, cw, ch, fill_color=RGBColor(0x18,0x1e,0x3a),
         line_color=RGBColor(0x30,0x38,0x60))
    txb(s, f"{icon}  {title}", x+Inches(0.15), y+Inches(0.12), cw-Inches(0.3), Inches(0.38),
        size=13, bold=True, color=GOLD)
    txb(s, desc, x+Inches(0.15), y+Inches(0.52), cw-Inches(0.3), ch-Inches(0.65),
        size=11, color=GRAY)

bottom_bar(s); slide_num(s, 2)

# ════════════════════════════════════════════════════════
# SLIDE 3 – Problem & Motivation
# ════════════════════════════════════════════════════════
s = add_slide(RGBColor(0x0a,0x0a,0x18))
section_label(s, "问题分析", RED)
txb(s, "现有方法的不足", Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=RED)

problems = [
    ("❌  问题一：特征一视同仁",
     "现有方法对所有通道和空间位置施以相同权重，忽略了不同特征对去雾贡献的差异性。"),
    ("❌  问题二：单一尺度特征",
     "大多数网络仅使用单一分支提取特征，无法同时捕获不同层次的去雾信息（全局/局部）。"),
    ("❌  问题三：融合策略简单",
     "多分支结果通常直接拼接或平均融合，未能自适应地学习各分支的重要性权重。"),
]
for i,(title,desc) in enumerate(problems):
    y = Inches(1.55) + i*Inches(1.4)
    rect(s, Inches(0.5), y, Inches(12.3), Inches(1.22),
         fill_color=RGBColor(0x25,0x0a,0x0a), line_color=RED)
    # left accent
    rect(s, Inches(0.5), y, Inches(0.07), Inches(1.22), fill_color=RED)
    txb(s, title, Inches(0.72), y+Inches(0.12), Inches(11.5), Inches(0.38),
        size=14, bold=True, color=RED)
    txb(s, desc, Inches(0.72), y+Inches(0.52), Inches(11.5), Inches(0.62),
        size=12, color=GRAY)

# insight box
rect(s, Inches(0.5), Inches(5.75), Inches(12.3), Inches(1.05),
     fill_color=RGBColor(0x05,0x18,0x2a), line_color=CYAN)
txb(s, "💡  核心洞察：设计双重注意力机制（通道 + 空间）并配合多分支特征融合注意力，"
       "让网络精准聚焦有效信息，大幅提升复原质量。",
    Inches(0.65), Inches(5.82), Inches(12.0), Inches(0.9),
    size=13, color=WHITE)

bottom_bar(s); slide_num(s, 3)

# ════════════════════════════════════════════════════════
# SLIDE 4 – Architecture Overview
# ════════════════════════════════════════════════════════
s = add_slide(BG_BLUE)
section_label(s, "网络架构", CYAN)
txb(s, "FFA-Net 整体架构", Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=CYAN)

# flow boxes
boxes = [
    ("输入\n有雾图像\n3×H×W",     WHITE),
    ("预处理层\nConv 3→64\n3×3",   CYAN),
    ("Group 1\nN×残差块",          GOLD),
    ("Group 2\nN×残差块",          GOLD),
    ("Group 3\nN×残差块",          GOLD),
    ("特征融合\n注意力加权\nConcat+CA", ORANGE),
    ("后处理层\nConv×2\n+残差",    CYAN),
    ("输出\n清晰图像\n3×H×W",      GREEN),
]

# draw group bracket
bx, by = Inches(3.9), Inches(1.52)
bw, bh = Inches(4.5), Inches(2.2)
rect(s, bx, by, bw, bh, fill_color=RGBColor(0x18,0x22,0x18),
     line_color=RGBColor(0x40,0x60,0x30))
txb(s, "并行 3 组", bx+bw/2-Inches(0.5), by-Inches(0.28), Inches(1.2), Inches(0.28),
    size=10, color=GOLD, align=PP_ALIGN.CENTER)

# flow
bxs = [Inches(0.18), Inches(2.15),
       Inches(4.05), Inches(5.58), Inches(7.12),
       Inches(8.65), Inches(10.35), Inches(11.7)]
bwy, bwx = Inches(1.05), Inches(1.4)
flow_y = Inches(2.1)
for i,(txt,clr) in enumerate(boxes):
    if i in (2,3,4):
        fy = flow_y + (i-2)*Inches(0.65) - Inches(0.1)
        fh = Inches(0.58)
    else:
        fy = flow_y + Inches(0.3)
        fh = bwy
    rect(s, bxs[i], fy, bwx, fh,
         fill_color=RGBColor(0x10,0x20,0x30), line_color=clr)
    txb(s, txt, bxs[i]+Inches(0.05), fy+Inches(0.04), bwx-Inches(0.1), fh-Inches(0.1),
        size=9, color=clr, align=PP_ALIGN.CENTER)
    # arrow
    if i < len(boxes)-1 and i not in (1,4):
        ax = bxs[i]+bwx+Inches(0.02)
        ay = fy + fh/2 - Inches(0.12)
        txb(s, "→", ax, ay, Inches(0.25), Inches(0.25),
            size=14, color=DGRAY, align=PP_ALIGN.CENTER)

txb(s, "全局跳跃连接：输出 = 网络预测 + 输入有雾图像",
    Inches(0.5), Inches(3.62), Inches(12), Inches(0.38),
    size=12, color=DGRAY, align=PP_ALIGN.CENTER, italic=True)

# three highlights
hi = [
    ("残差块设计", CYAN, "每个残差块含 2 个 3×3 卷积 + CA + PA 双重注意力"),
    ("并行多分支", GOLD, "3 组独立残差组并行，各含 19 个残差块，捕获互补信息"),
    ("自适应融合", RED,  "三组特征通过注意力学习权重 w₁、w₂、w₃，加权求和"),
]
for i,(title,clr,desc) in enumerate(hi):
    x = Inches(0.5) + i*Inches(4.28)
    rect(s, x, Inches(4.1), Inches(4.0), Inches(1.18),
         fill_color=RGBColor(0x10,0x18,0x28), line_color=clr)
    rect(s, x, Inches(4.1), Inches(4.0), Inches(0.07), fill_color=clr)
    txb(s, title, x+Inches(0.12), Inches(4.22), Inches(3.8), Inches(0.38),
        size=13, bold=True, color=clr)
    txb(s, desc, x+Inches(0.12), Inches(4.62), Inches(3.8), Inches(0.58),
        size=11, color=GRAY)

bottom_bar(s); slide_num(s, 4)

# ════════════════════════════════════════════════════════
# SLIDE 5 – Dual Attention
# ════════════════════════════════════════════════════════
s = add_slide(RGBColor(0x0f,0x20,0x27))
section_label(s, "关键模块", GREEN)
txb(s, "双重注意力机制", Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=GREEN)

for col,(title,clr,steps) in enumerate([
    ("🟢  通道注意力 (Channel Attention, CA)", GREEN, [
        "① 输入特征图 F（C × H × W）",
        "② 全局平均池化 → (C × 1 × 1)  压缩空间信息",
        "③ 两层 1×1 卷积 + ReLU / Sigmoid",
        "④ 通道权重向量与 F 逐通道相乘（重新校准）",
        "⑤ 输出：增强了重要通道的特征图",
        "→ 让网络关注哪些通道对去雾最有效",
    ]),
    ("🟠  空间注意力 (Pixel Attention, PA)", ORANGE, [
        "① 输入特征图 F（C × H × W）",
        "② 两层 1×1 卷积将通道数压缩至 1",
        "③ Sigmoid 激活 → 空间掩码 M（1 × H × W）",
        "④ 掩码 M 与 F 逐像素相乘（空间重加权）",
        "⑤ 输出：突出雾浓区域的特征图",
        "→ 让网络关注哪些空间位置的雾更浓",
    ]),
]):
    x = Inches(0.5) + col*Inches(6.45)
    rect(s, x, Inches(1.52), Inches(6.2), Inches(4.3),
         fill_color=RGBColor(0x10,0x18,0x22), line_color=clr)
    rect(s, x, Inches(1.52), Inches(6.2), Inches(0.06), fill_color=clr)
    txb(s, title, x+Inches(0.15), Inches(1.62), Inches(6.0), Inches(0.45),
        size=13, bold=True, color=clr)
    for j,step in enumerate(steps):
        sy = Inches(2.18) + j*Inches(0.55)
        fc = RGBColor(0x20,0x35,0x20) if col==0 else RGBColor(0x30,0x20,0x10)
        rect(s, x+Inches(0.15), sy, Inches(5.8), Inches(0.44),
             fill_color=fc, line_color=None)
        txb(s, step, x+Inches(0.25), sy+Inches(0.06), Inches(5.6), Inches(0.35),
            size=11, color=WHITE if j<5 else clr, bold=(j==5))

# formula box
rect(s, Inches(0.5), Inches(6.02), Inches(12.3), Inches(0.82),
     fill_color=RGBColor(0x10,0x20,0x18), line_color=GREEN)
txb(s, "残差块输出公式：  F_out  =  PA( CA( Conv₂( Conv₁(F) ) ) )  +  F",
    Inches(0.7), Inches(6.12), Inches(12), Inches(0.62),
    size=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

bottom_bar(s); slide_num(s, 5)

# ════════════════════════════════════════════════════════
# SLIDE 6 – Feature Fusion Attention
# ════════════════════════════════════════════════════════
s = add_slide(RGBColor(0x14,0x05,0x28))
section_label(s, "关键模块", PURPLE)
txb(s, "特征融合注意力（FFA 模块）", Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=PURPLE)

# diagram boxes
groups = [("Group 1 输出", CYAN), ("Group 2 输出", GOLD), ("Group 3 输出", RED)]
for i,(lbl,clr) in enumerate(groups):
    y = Inches(1.6) + i*Inches(1.05)
    rect(s, Inches(0.5), y, Inches(2.4), Inches(0.82),
         fill_color=RGBColor(0x10,0x18,0x28), line_color=clr)
    txb(s, lbl, Inches(0.55), y+Inches(0.22), Inches(2.3), Inches(0.42),
        size=12, color=clr, align=PP_ALIGN.CENTER)

# arrows
for i in range(3):
    y = Inches(1.95) + i*Inches(1.05)
    txb(s, "→", Inches(2.95), y-Inches(0.06), Inches(0.4), Inches(0.35),
        size=16, color=DGRAY, align=PP_ALIGN.CENTER)

# concat
rect(s, Inches(3.4), Inches(1.55), Inches(1.9), Inches(3.1),
     fill_color=RGBColor(0x20,0x10,0x35), line_color=PURPLE)
txb(s, "Concat\n拼接\n(3C×H×W)", Inches(3.42), Inches(2.4), Inches(1.86), Inches(1.2),
    size=11, color=PURPLE, align=PP_ALIGN.CENTER)

txb(s, "→", Inches(5.35), Inches(2.9), Inches(0.4), Inches(0.35),
    size=16, color=DGRAY, align=PP_ALIGN.CENTER)

# attention
rect(s, Inches(5.8), Inches(1.95), Inches(2.2), Inches(2.3),
     fill_color=RGBColor(0x28,0x18,0x08), line_color=ORANGE)
txb(s, "通道注意力\n学习权重\nw₁, w₂, w₃", Inches(5.82), Inches(2.5), Inches(2.16), Inches(1.2),
    size=11, color=ORANGE, align=PP_ALIGN.CENTER)

txb(s, "→", Inches(8.05), Inches(2.9), Inches(0.4), Inches(0.35),
    size=16, color=DGRAY, align=PP_ALIGN.CENTER)

# output
rect(s, Inches(8.5), Inches(2.2), Inches(2.3), Inches(1.82),
     fill_color=RGBColor(0x08,0x25,0x15), line_color=GREEN)
txb(s, "融合输出\n加权求和", Inches(8.52), Inches(2.7), Inches(2.26), Inches(1.0),
    size=12, color=GREEN, align=PP_ALIGN.CENTER)

# formula
rect(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.62),
     fill_color=RGBColor(0x20,0x10,0x38), line_color=PURPLE)
txb(s, "Out = w₁·F₁  +  w₂·F₂  +  w₃·F₃     （权重由网络自动学习，w₁+w₂+w₃=1）",
    Inches(0.65), Inches(5.0), Inches(12.0), Inches(0.45),
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# 4 desc items
descs = [
    ("自适应权重", "权重由网络自动学习，而非人工固定，可根据输入内容动态调整各分支贡献"),
    ("互补信息",   "三个并行分支各自独立，产生互补去雾特征，融合后信息更完整"),
    ("端到端优化", "融合权重与网络参数联合训练，L1 损失 + 可选 VGG16 感知损失统一优化"),
    ("命名由来",   "FFA = Feature Fusion Attention，此模块正是论文核心贡献"),
]
for i,(t,d) in enumerate(descs):
    x = Inches(0.5) + i*Inches(3.22)
    rect(s, x, Inches(5.68), Inches(3.05), Inches(1.12),
         fill_color=RGBColor(0x15,0x08,0x25), line_color=PURPLE)
    txb(s, t, x+Inches(0.12), Inches(5.73), Inches(2.85), Inches(0.35),
        size=12, bold=True, color=PURPLE)
    txb(s, d, x+Inches(0.12), Inches(6.1), Inches(2.85), Inches(0.6),
        size=10, color=GRAY)

bottom_bar(s); slide_num(s, 6)

# ════════════════════════════════════════════════════════
# SLIDE 7 – Dataset & Training
# ════════════════════════════════════════════════════════
s = add_slide(RGBColor(0x08,0x20,0x18))
section_label(s, "实验设置", GREEN)
txb(s, "数据集：RESIDE Benchmark", Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=GREEN)

cards_data = [
    ("室内场景", [
        ("训练集 (ITS)",     "13,990 对"),
        ("测试集 (SOTS)",    "500 对"),
        ("图像格式",         "PNG"),
        ("雾浓度 (β)",       "0.6 – 1.8"),
        ("合成方式",         "深度图 + 散射模型"),
    ]),
    ("室外场景", [
        ("训练集 (OTS)",     "313,950 对"),
        ("测试集 (SOTS)",    "500 对"),
        ("图像格式",         "JPEG"),
        ("雾浓度",           "多种浓度混合"),
        ("合成方式",         "真实场景 + 大气散射"),
    ]),
    ("训练超参数", [
        ("残差块数 (blocks)", "19"),
        ("分支数 (gps)",      "3"),
        ("批大小 (bs)",       "2"),
        ("学习率 (lr)",       "1e-4"),
        ("裁剪尺寸",          "240 × 240"),
    ]),
    ("训练策略", [
        ("室内训练步数",   "500,000"),
        ("室外训练步数",   "1,000,000"),
        ("学习率调度",     "余弦退火"),
        ("数据增强",       "随机翻转 + 旋转"),
        ("评估指标",       "PSNR / SSIM"),
    ]),
]
for i,(title,items) in enumerate(cards_data):
    c,r = i%2, i//2
    x = Inches(0.5) + c*Inches(6.45)
    y = Inches(1.55) + r*Inches(2.4)
    rect(s, x, y, Inches(6.2), Inches(2.2),
         fill_color=RGBColor(0x10,0x22,0x18), line_color=GREEN)
    txb(s, title, x+Inches(0.15), y+Inches(0.1), Inches(6.0), Inches(0.38),
        size=13, bold=True, color=GREEN)
    for j,(k,v) in enumerate(items):
        iy = y + Inches(0.55) + j*Inches(0.3)
        txb(s, k, x+Inches(0.15), iy, Inches(3.0), Inches(0.28), size=11, color=GRAY)
        txb(s, v, x+Inches(3.2), iy, Inches(2.8), Inches(0.28), size=11, color=GOLD, bold=True)

rect(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.58),
     fill_color=RGBColor(0x08,0x28,0x18), line_color=GREEN)
txb(s, "📌  损失函数：L₁ 重建损失（必选）+ VGG16 感知损失（可选），两者均支持余弦退火学习率调度",
    Inches(0.65), Inches(6.57), Inches(12.0), Inches(0.45), size=12, color=WHITE)

bottom_bar(s); slide_num(s, 7)

# ════════════════════════════════════════════════════════
# SLIDE 8 – Quantitative Results
# ════════════════════════════════════════════════════════
s = add_slide(RGBColor(0x0a,0x0a,0x18))
section_label(s, "实验结果", GOLD)
txb(s, "定量对比（SOTS 数据集，PSNR / SSIM）", Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=26, bold=True, color=GOLD)

headers = ["方法", "年份", "类型", "室内 PSNR↑", "室内 SSIM↑", "室外 PSNR↑", "室外 SSIM↑"]
rows_data = [
    ("DCP",         "2009", "传统",    "16.62", "0.8179", "19.13", "0.8148", False),
    ("AOD-Net",     "2017", "CNN",     "19.06", "0.8504", "20.29", "0.8765", False),
    ("DehazeNet",   "2016", "CNN",     "21.14", "0.8472", "22.46", "0.8514", False),
    ("GFN",         "2018", "CNN",     "22.30", "0.8800", "21.55", "0.8444", False),
    ("GCANet",      "2019", "CNN+Attn","30.23", "0.9800", "—",     "—",      False),
    ("FFA-Net (Ours)","2020","双重注意力","36.39","0.9886","33.57","0.9840", True),
]
col_ws = [Inches(2.0), Inches(0.8), Inches(1.4), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)]
col_xs = [Inches(0.5)]
for w in col_ws[:-1]:
    col_xs.append(col_xs[-1]+w+Inches(0.05))

# header row
rect(s, Inches(0.5), Inches(1.55), sum(col_ws)+Inches(0.3), Inches(0.48),
     fill_color=RGBColor(0x28,0x22,0x05), line_color=GOLD)
for i,h in enumerate(headers):
    txb(s, h, col_xs[i]+Inches(0.05), Inches(1.6), col_ws[i]-Inches(0.05), Inches(0.38),
        size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

for ri,(row) in enumerate(rows_data):
    *cells, is_best = row
    ry = Inches(2.1) + ri*Inches(0.68)
    bg = RGBColor(0x22,0x1a,0x04) if is_best else RGBColor(0x12,0x14,0x22)
    lc = GOLD if is_best else RGBColor(0x28,0x28,0x40)
    rect(s, Inches(0.5), ry, sum(col_ws)+Inches(0.3), Inches(0.6),
         fill_color=bg, line_color=lc)
    for ci,cell in enumerate(cells):
        clr = GOLD if (is_best and ci>=3) else (WHITE if is_best else GRAY)
        txb(s, cell, col_xs[ci]+Inches(0.05), ry+Inches(0.12),
            col_ws[ci]-Inches(0.05), Inches(0.38),
            size=12, color=clr, bold=is_best,
            align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)

# notes
notes = [
    ("室内提升", "vs GCANet：PSNR +6.16 dB\nvs DehazeNet：+15.25 dB"),
    ("室外提升", "vs GFN：PSNR +12.02 dB\nvs DehazeNet：+11.11 dB"),
    ("SSIM 指标", "室内 0.9886 / 室外 0.9840\n接近满分（1.0）"),
]
for i,(t,d) in enumerate(notes):
    x = Inches(0.5) + i*Inches(4.28)
    rect(s, x, Inches(6.7), Inches(4.0), Inches(0.95),
         fill_color=RGBColor(0x20,0x18,0x05), line_color=GOLD)
    txb(s, t, x+Inches(0.12), Inches(6.75), Inches(3.8), Inches(0.35),
        size=12, bold=True, color=GOLD)
    txb(s, d, x+Inches(0.12), Inches(7.1), Inches(3.8), Inches(0.5),
        size=10, color=GRAY)

bottom_bar(s); slide_num(s, 8)

# ════════════════════════════════════════════════════════
# SLIDE 9 – Visual Results
# ════════════════════════════════════════════════════════
s = add_slide(BG_BLACK)
section_label(s, "实验结果", CYAN)
txb(s, "定性对比（视觉效果展示）", Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=CYAN)

fig_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "fig")

pairs = [
    ("室外场景 — 建筑物",
     os.path.join(fig_dir,"1400_2.png"),
     os.path.join(fig_dir,"1400_2_FFA.png")),
    ("室内场景 — 居室",
     os.path.join(fig_dir,"0099_0.9_0.16.jpg"),
     os.path.join(fig_dir,"0099_0_FFA.png")),
]

for col,(label, hazy_path, clear_path) in enumerate(pairs):
    ox = Inches(0.5) + col*Inches(6.5)
    txb(s, label, ox, Inches(1.52), Inches(6.2), Inches(0.38),
        size=13, color=GRAY)
    for j,(path, lbl) in enumerate([(hazy_path,"输入（有雾）"),(clear_path,"FFA-Net 输出")]):
        ix = ox + j*Inches(3.05)
        rect(s, ix, Inches(1.98), Inches(2.85), Inches(3.6),
             fill_color=RGBColor(0x18,0x18,0x18), line_color=RGBColor(0x33,0x33,0x33))
        if os.path.exists(path):
            try:
                s.shapes.add_picture(path, ix+Inches(0.05), Inches(2.0),
                                     Inches(2.75), Inches(3.0))
            except Exception:
                txb(s, "[图片]", ix+Inches(0.8), Inches(3.2), Inches(1.5), Inches(0.5),
                    size=11, color=DGRAY, align=PP_ALIGN.CENTER)
        txb(s, lbl, ix, Inches(5.05), Inches(2.85), Inches(0.38),
            size=11, color=DGRAY, align=PP_ALIGN.CENTER)

# obs notes
obs = [
    ("颜色保真", "去雾后色彩自然，无过饱和或色偏"),
    ("细节恢复", "远景轮廓清晰，纹理细节完整保留"),
    ("对比度",   "整体对比度提升，层次感显著增强"),
]
for i,(t,d) in enumerate(obs):
    x = Inches(0.5) + i*Inches(4.28)
    rect(s, x, Inches(5.65), Inches(4.0), Inches(0.95),
         fill_color=RGBColor(0x05,0x18,0x28), line_color=CYAN)
    txb(s, t, x+Inches(0.12), Inches(5.7), Inches(3.8), Inches(0.35),
        size=12, bold=True, color=CYAN)
    txb(s, d, x+Inches(0.12), Inches(6.05), Inches(3.8), Inches(0.48),
        size=11, color=GRAY)

bottom_bar(s); slide_num(s, 9)

# ════════════════════════════════════════════════════════
# SLIDE 10 – Ablation Study
# ════════════════════════════════════════════════════════
s = add_slide(RGBColor(0x10,0x08,0x22))
section_label(s, "消融实验", PURPLE)
txb(s, "消融实验分析", Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=PURPLE)

abl_cols = ["变体", "CA", "PA", "FFA", "多分支", "室内 PSNR", "室外 PSNR"]
abl_ws   = [Inches(3.0), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9), Inches(1.5), Inches(1.5)]
abl_xs   = [Inches(0.5)]
for w in abl_ws[:-1]:
    abl_xs.append(abl_xs[-1]+w+Inches(0.05))

rect(s, Inches(0.5), Inches(1.52), sum(abl_ws)+Inches(0.3), Inches(0.45),
     fill_color=RGBColor(0x25,0x10,0x42), line_color=PURPLE)
for i,h in enumerate(abl_cols):
    txb(s, h, abl_xs[i]+Inches(0.05), Inches(1.56), abl_ws[i]-Inches(0.05), Inches(0.36),
        size=12, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

abl_rows = [
    ("基础 CNN（无注意力）", "✗","✗","✗","✗","30.41","26.87", False),
    ("+ 通道注意力 (CA)",    "✓","✗","✗","✗","32.15","28.36", False),
    ("+ 空间注意力 (PA)",    "✗","✓","✗","✗","31.78","27.94", False),
    ("CA + PA（双重）",      "✓","✓","✗","✗","33.82","30.11", False),
    ("CA + PA + 多分支",     "✓","✓","✗","✓","35.21","32.08", False),
    ("完整 FFA-Net",         "✓","✓","✓","✓","36.39","33.57", True),
]
for ri,row in enumerate(abl_rows):
    *cells, is_best = row
    ry = Inches(2.05) + ri*Inches(0.65)
    bg = RGBColor(0x22,0x10,0x38) if is_best else RGBColor(0x14,0x10,0x24)
    lc = PURPLE if is_best else RGBColor(0x28,0x20,0x40)
    rect(s, Inches(0.5), ry, sum(abl_ws)+Inches(0.3), Inches(0.58),
         fill_color=bg, line_color=lc)
    for ci,cell in enumerate(cells):
        if ci == 0:
            clr = WHITE if is_best else GRAY
            aln = PP_ALIGN.LEFT
        elif cell == "✓":
            clr = GREEN; aln = PP_ALIGN.CENTER
        elif cell == "✗":
            clr = RGBColor(0x44,0x44,0x44); aln = PP_ALIGN.CENTER
        else:
            clr = GOLD if is_best else RGBColor(0xFF,0xD2,0x80)
            aln = PP_ALIGN.CENTER
        txb(s, cell, abl_xs[ci]+Inches(0.05), ry+Inches(0.1),
            abl_ws[ci]-Inches(0.05), Inches(0.38),
            size=12, color=clr, bold=is_best, align=aln)

rect(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.95),
     fill_color=RGBColor(0x15,0x08,0x28), line_color=PURPLE)
txb(s, "结论：CA 和 PA 各自有效，合并更好；多分支带来约 +1.4 dB；"
       "FFA 融合注意力再带来约 +1.2 dB；三者协同达到最佳性能 36.39 dB。",
    Inches(0.65), Inches(6.42), Inches(12.0), Inches(0.8),
    size=12, color=WHITE)

bottom_bar(s); slide_num(s, 10)

# ════════════════════════════════════════════════════════
# SLIDE 11 – Complexity
# ════════════════════════════════════════════════════════
s = add_slide(RGBColor(0x08,0x12,0x22))
section_label(s, "模型分析", CYAN)
txb(s, "模型复杂度与性能权衡", Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=CYAN)

# left: bar chart
rect(s, Inches(0.5), Inches(1.52), Inches(6.2), Inches(4.2),
     fill_color=RGBColor(0x10,0x18,0x28), line_color=RGBColor(0x28,0x38,0x58))
txb(s, "各方法参数量对比（M）", Inches(0.65), Inches(1.6), Inches(6.0), Inches(0.38),
    size=13, bold=True, color=CYAN)

bars = [
    ("FFA-Net", 4.68, GOLD),
    ("GCANet",  0.70, CYAN),
    ("GFN",     0.58, GREEN),
    ("DehazeNet",0.008,RED),
    ("AOD-Net", 0.002, PURPLE),
]
max_val = 5.0
bar_area_w = Inches(4.5)
for bi,(name,val,clr) in enumerate(bars):
    by = Inches(2.1) + bi*Inches(0.62)
    txb(s, f"{name}", Inches(0.65), by, Inches(1.5), Inches(0.3),
        size=11, color=GRAY)
    bw = bar_area_w * (val / max_val)
    rect(s, Inches(2.2), by+Inches(0.04), max(bw, Inches(0.1)), Inches(0.25),
         fill_color=clr, line_color=None)
    txb(s, f"{val} M", Inches(2.2)+bw+Inches(0.05), by, Inches(0.8), Inches(0.3),
        size=10, color=clr)

# right: insight
rect(s, Inches(6.85), Inches(1.52), Inches(6.0), Inches(4.2),
     fill_color=RGBColor(0x10,0x18,0x28), line_color=RGBColor(0x28,0x38,0x58))
txb(s, "性能 vs 参数量分析", Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.38),
    size=13, bold=True, color=CYAN)

insights = [
    ("参数增加", "~4.68 M vs GCANet 0.70 M", "约 6.7 倍"),
    ("PSNR 增益", "36.39 vs 30.23（GCANet）", "+6.16 dB"),
    ("SSIM 增益", "0.9886 vs 0.9800（GCANet）", "+0.0086"),
]
for ii,(t,sub,val) in enumerate(insights):
    iy = Inches(2.1) + ii*Inches(1.0)
    rect(s, Inches(7.0), iy, Inches(5.5), Inches(0.85),
         fill_color=RGBColor(0x08,0x18,0x30), line_color=RGBColor(0x18,0x38,0x60))
    txb(s, t, Inches(7.12), iy+Inches(0.08), Inches(2.5), Inches(0.32),
        size=12, bold=True, color=CYAN)
    txb(s, sub, Inches(7.12), iy+Inches(0.42), Inches(3.5), Inches(0.32),
        size=10, color=GRAY)
    txb(s, val, Inches(10.0), iy+Inches(0.22), Inches(2.3), Inches(0.42),
        size=14, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

rect(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.85),
     fill_color=RGBColor(0x05,0x15,0x28), line_color=CYAN)
txb(s, "总结：以约 6.7 倍参数量换取 6.16 dB PSNR 增益，性价比显著。"
       "若资源受限可减小 blocks 数（如设为 6）获得轻量变体；"
       "推理支持全分辨率输入，单张 GPU 约 0.1~0.5 秒。",
    Inches(0.65), Inches(5.97), Inches(12.0), Inches(0.72),
    size=12, color=WHITE)

bottom_bar(s); slide_num(s, 11)

# ════════════════════════════════════════════════════════
# SLIDE 12 – Conclusion
# ════════════════════════════════════════════════════════
s = add_slide(BG_DARK)
section_label(s, "总结", GOLD)
txb(s, "结论与展望", Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=GOLD)

con_items = [
    ("01", "核心贡献",
     "提出双重注意力（CA+PA）与特征融合注意力（FFA）相结合的去雾网络，"
     "在 SOTS 数据集上大幅超越所有已知方法，室内 PSNR 达 36.39 dB。"),
    ("02", "方法创新",
     "首次在去雾任务中同时使用通道和空间两维注意力，并设计自适应多分支特征融合，"
     "摆脱了手工估算透射率的局限。"),
    ("03", "实用价值",
     "代码开源，提供完整训练/测试流程和预训练模型，支持室内/室外两场景，"
     "可直接应用于工程实践。"),
    ("04", "启示影响",
     "双重注意力设计对超分辨率、图像增强、低光照复原等任务具有普遍参考价值，"
     "已被后续多篇工作引用。"),
]
for i,(num,title,desc) in enumerate(con_items):
    c, r = i%2, i//2
    x = Inches(0.5)  + c*Inches(6.45)
    y = Inches(1.55) + r*Inches(2.35)
    rect(s, x, y, Inches(6.2), Inches(2.15),
         fill_color=RGBColor(0x18,0x15,0x30), line_color=RGBColor(0x44,0x35,0x88))
    txb(s, num, x+Inches(0.12), y+Inches(0.08), Inches(0.8), Inches(0.65),
        size=32, bold=True, color=RGBColor(0x33,0x28,0x55))
    txb(s, title, x+Inches(0.12), y+Inches(0.12), Inches(6.0), Inches(0.4),
        size=14, bold=True, color=GOLD)
    txb(s, desc, x+Inches(0.12), y+Inches(0.58), Inches(5.95), Inches(1.5),
        size=11, color=GRAY)

rect(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.88),
     fill_color=RGBColor(0x25,0x0a,0x0a), line_color=RED)
txb(s, "局限与未来方向：  ① 在真实（非合成）有雾图像上的泛化性有待验证；"
       "  ② 参数量较大，移动端需要模型压缩；"
       "  ③ 未来可探索 Transformer 自注意力与本文卷积注意力的结合。",
    Inches(0.65), Inches(6.47), Inches(12.0), Inches(0.75),
    size=11, color=WHITE)

bottom_bar(s); slide_num(s, 12)

# ════════════════════════════════════════════════════════
# SLIDE 13 – Thank You
# ════════════════════════════════════════════════════════
s = add_slide(BG_DARK)
txb(s, "谢谢！", Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.6),
    size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txb(s, "欢迎提问与讨论",
    Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.65),
    size=24, color=GRAY, align=PP_ALIGN.CENTER)

divider(s, Inches(4.35))

txb(s, "参考文献", Inches(2.0), Inches(4.55), Inches(9.3), Inches(0.42),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s,
    "Qin X, Wang Z, Bai Y, et al. FFA-Net: Feature fusion attention network for single image dehazing[C].\n"
    "Proceedings of the AAAI Conference on Artificial Intelligence. 2020, 34(07): 11908–11915.\n"
    "arXiv: 1911.07559  |  数据集: RESIDE Benchmark  |  代码: GitHub (Bryant1iu/FFA-Net)",
    Inches(1.5), Inches(5.05), Inches(10.3), Inches(1.5),
    size=12, color=GRAY, align=PP_ALIGN.CENTER)

bottom_bar(s); slide_num(s, 13)

# ── Save ──────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "FFA-Net_presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
